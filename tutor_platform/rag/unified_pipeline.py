"""Unified document processing pipeline — single entry point for all ingestion paths.

Routes every file through classification → extraction → optional structuring,
producing standardized sidecar files regardless of the original entry point
(Web KB upload, WeChat, MCP tools, CLI, chat attachments, API calls).

Usage::

    result = await UnifiedDocumentPipeline.process("/path/to/file.pdf")
    # → produces .txt and/or .exam.json sidecar next to the original file
"""

from __future__ import annotations

import asyncio
import logging
import os
from enum import Enum
from pathlib import Path
from typing import TYPE_CHECKING

logger = logging.getLogger(__name__)

if TYPE_CHECKING:
    from deeptutor.services.llm.client import LLMClient


# ── File classification ─────────────────────────────────────────

class DocType(str, Enum):
    """Document type determined by extension + content sniffing."""

    TEXT = "text"               # .txt, .md, .csv, .json, etc.
    IMAGE = "image"             # .jpg, .png, .webp, etc.
    TEXT_PDF = "text_pdf"        # PDF with extractable text layer
    SCANNED_PDF = "scanned_pdf"   # PDF without text layer
    EXAM_PDF = "exam_pdf"        # PDF that looks like an exam (multi-page, contains questions)
    OFFICE_DOCX = "office_docx"   # .docx
    OFFICE_XLSX = "office_xlsx"   # .xlsx, .xls
    OFFICE_PPTX = "office_pptx"   # .pptx, .ppt
    OFFICE_OLD = "office_old"     # .doc, .ppt, .xls
    OFFICE_OTHER = "office_other"  # .odt, .rtf, .pdf (text), etc.
    UNKNOWN = "unknown"


# Extension → DocType mapping (fast path, no content sniffing)
_EXT_TYPE_MAP: dict[str, DocType] = {
    ".txt": DocType.TEXT, ".md": DocType.TEXT, ".csv": DocType.TEXT,
    ".json": DocType.TEXT, ".yaml": DocType.TEXT, ".yml": DocType.TEXT,
    ".xml": DocType.TEXT, ".html": DocType.TEXT, ".htm": DocType.TEXT,
    ".py": DocType.TEXT, ".js": DocType.TEXT, ".ts": DocType.TEXT,
    ".css": DocType.TEXT, ".log": DocType.TEXT,
    ".jpg": DocType.IMAGE, ".jpeg": DocType.IMAGE, ".png": DocType.IMAGE,
    ".gif": DocType.IMAGE, ".webp": DocType.IMAGE, ".bmp": DocType.IMAGE,
    ".tiff": DocType.IMAGE, ".tif": DocType.IMAGE, ".heic": DocType.IMAGE,
    ".docx": DocType.OFFICE_DOCX,
    ".doc": DocType.OFFICE_OLD,
    ".xlsx": DocType.OFFICE_XLSX,
    ".xls": DocType.OFFICE_OLD,
    ".pptx": DocType.OFFICE_PPTX, ".pptm": DocType.OFFICE_PPTX,
    ".ppt": DocType.OFFICE_OLD, ".pps": DocType.OFFICE_OLD, ".ppsx": DocType.OFFICE_PPTX,
    ".odt": DocType.OFFICE_OTHER, ".rtf": DocType.OFFICE_OTHER,
}


def classify_file(file_path: str | Path) -> DocType:
    """Classify a file by extension and content sniffing.

    PDFs get further classified into text_pdf / scanned_pdf / exam_pdf.
    """
    ext = os.path.splitext(str(file_path))[1].lower()
    doc_type = _EXT_TYPE_MAP.get(ext, DocType.UNKNOWN)

    if doc_type == DocType.UNKNOWN and ext == ".pdf":
        return _classify_pdf(file_path)

    return doc_type


def _classify_pdf(file_path: str | Path) -> DocType:
    """Open a PDF with pymupdf and decide text-layer vs scanned."""
    try:
        import fitz
    except ImportError:
        return DocType.SCANNED_PDF

    path = Path(file_path)
    try:
        doc = fitz.open(path)
    except Exception:
        return DocType.SCANNED_PDF

    try:
        total_chars = 0
        image_blocks = 0
        for page in doc:
            text = page.get_text()
            total_chars += len(text.strip())

            blocks = page.get_text("dict").get("blocks", [])
            image_blocks += sum(1 for b in blocks if b.get("type") == 1)

        # Heuristics for exam PDF:
        # - Multi-page (>2) with image blocks (questions with figures)
        # - OR single-page with both text AND at least one image block
        is_exam = (len(doc) > 2 and image_blocks > 0) or (image_blocks > 0 and total_chars > 200)

        if is_exam:
            return DocType.EXAM_PDF

        if total_chars > 100:
            return DocType.TEXT_PDF
        else:
            return DocType.SCANNED_PDF
    finally:
        doc.close()


# ── Unified pipeline result ──────────────────────────────────────

class PipelineResult:
    """Result of running the unified pipeline on a single file."""

    __slots__ = (
        "file_path", "doc_type", "content_text", "sidecar_paths",
        "error", "stats",
    )

    def __init__(
        self,
        file_path: str,
        doc_type: DocType,
        content_text: str | None = None,
        sidecar_paths: list[str] | None = None,
        error: str | None = None,
        stats: dict | None = None,
    ) -> None:
        self.file_path = file_path
        self.doc_type = doc_type
        self.content_text = content_text
        self.sidecar_paths = sidecar_paths or []
        self.error = error
        self.stats = stats or {}

    @property
    def ok(self) -> bool:
        return self.error is None


# ── Unified pipeline ─────────────────────────────────────────────

class UnifiedDocumentPipeline:
    """Process any document through classification → extraction → structuring.

    Produces standardized sidecar files alongside the original:

    - ``.txt`` — plain-text extraction (all document types)
    - ``.exam.json`` — structured exam paper (exam PDFs only)
    - ``.figures/`` — clipped figure PNGs (exam PDFs with images)

    Usage::

        result = await UnifiedDocumentPipeline.process("/tmp/exam.pdf")
        print(result.content_text)       # extracted plain text
        print(result.sidecar_paths)       # [".txt", ".exam.json"]
    """

    @classmethod
    async def process(
        cls,
        file_path: str | Path,
        *,
        llm_client: "LLMClient | None" = None,
        enable_structured_exam: bool = True,
        max_exam_pages: int = 50,
    ) -> PipelineResult:
        """Run the unified pipeline on a single file.

        Args:
            file_path: Path to the file to process.
            llm_client: Optional LLMClient for multimodal OCR/vision.
                        If None, text extraction is best-effort without LLM.
            enable_structured_exam: When True, run Phase 1-4 for exam PDFs.
            max_exam_pages: Maximum pages to process for exam PDFs.
        """
        path = Path(file_path)
        stats: dict = {"file_type": None, "extraction_chars": 0, "sidecars": []}

        if not path.is_file():
            return PipelineResult(
                file_path=str(path),
                doc_type=DocType.UNKNOWN,
                error=f"File not found: {file_path}",
            )

        # 1. Classify
        doc_type = classify_file(path)
        stats["file_type"] = doc_type.value
        logger.info("Unified pipeline: %s → %s", path.name, doc_type.value)

        # 2. Extract text
        content_text: str | None = None
        sidecar_paths: list[str] = []

        try:
            if doc_type in (DocType.TEXT_PDF, DocType.SCANNED_PDF, DocType.EXAM_PDF):
                content_text = await cls._extract_pdf_text(path, doc_type, llm_client, stats)
            elif doc_type == DocType.IMAGE:
                content_text = await cls._extract_image_text(path, llm_client, stats)
            elif doc_type in (
                DocType.OFFICE_DOCX, DocType.OFFICE_XLSX, DocType.OFFICE_PPTX,
                DocType.OFFICE_OLD, DocType.OFFICE_OTHER,
            ):
                content_text = cls._extract_office_text(path, stats)
            elif doc_type == DocType.TEXT:
                content_text = cls._extract_text_file(path, stats)
            else:
                content_text = cls._extract_unknown(path, stats)

            stats["extraction_chars"] = len(content_text) if content_text else 0
        except Exception as exc:
            logger.error("Extraction failed for %s: %s", path.name, exc)
            return PipelineResult(
                file_path=str(path), doc_type=doc_type,
                error=f"Extraction failed: {exc}", stats=stats,
            )

        if not content_text:
            return PipelineResult(
                file_path=str(path), doc_type=doc_type,
                content_text="", stats=stats,
            )

        # 3. Write plain-text sidecar
        txt_path = cls._write_sidecar(path, content_text, suffix=".txt")
        if txt_path:
            sidecar_paths.append(txt_path)
            stats["sidecars"].append(".txt")

        # 4. Structured exam pipeline (Phase 1-4) for exam PDFs
        if enable_structured_exam and doc_type == DocType.EXAM_PDF and llm_client:
            try:
                exam_json = await cls._run_exam_pipeline(
                    path, llm_client, max_exam_pages, stats,
                )
                if exam_json:
                    json_path = cls._write_sidecar_raw(path, exam_json, suffix=".exam.json")
                    if json_path:
                        sidecar_paths.append(json_path)
                        stats["sidecars"].append(".exam.json")
            except Exception as exc:
                logger.error("Exam pipeline failed for %s: %s", path.name, exc)
                stats["exam_pipeline_error"] = str(exc)

        return PipelineResult(
            file_path=str(path),
            doc_type=doc_type,
            content_text=content_text,
            sidecar_paths=sidecar_paths,
            stats=stats,
        )

    # ── Per-type extractors ──────────────────────────────────────

    @classmethod
    async def _extract_pdf_text(
        cls, path: Path, doc_type: DocType, llm_client, stats: dict,
    ) -> str | None:
        """Extract text from any PDF type."""
        # For text-layer PDFs, use pymupdf directly
        if doc_type == DocType.TEXT_PDF:
            return _extract_text_pdf_fast(path)

        # For scanned/exam PDFs: try pymupdf first, then OCR
        text = _extract_text_pdf_fast(path)
        if text and len(text.strip()) > 50:
            return text

        # Need OCR — requires LLM client
        if llm_client is None:
            logger.warning("No LLM client available for scanned PDF %s", path.name)
            return text or ""

        return await cls._extract_scanned_pdf_text(path, llm_client, stats)

    @classmethod
    async def _extract_scanned_pdf_text(
        cls, path: Path, llm_client, stats: dict,
    ) -> str:
        """OCR a scanned PDF using the multimodal LLM."""
        try:
            import fitz
        except ImportError:
            return ""

        doc = fitz.open(path)
        try:
            pages_text: list[str] = []
            for i in range(len(doc)):
                page = doc[i]
                pix = page.get_pixmap(dpi=200)
                img_b64 = _to_base64(pix.tobytes("png"))

                try:
                    result = await llm_client.complete(
                        "Transcribe all text from this document page. Return only the text.",
                        image_data=img_b64,
                        image_mime_type="image/png",
                        image_filename=f"{path.name}:page{i+1}",
                    )
                    pages_text.append(result.strip())
                except Exception as exc:
                    logger.warning("OCR page %d of %s failed: %s", i + 1, path.name, exc)

            stats["ocr_pages"] = len(pages_text)
            return "\n\n".join(pages_text)
        finally:
            doc.close()

    @classmethod
    async def _extract_image_text(
        cls, path: Path, llm_client, stats: dict,
    ) -> str | None:
        """OCR an image file."""
        if llm_client is None:
            return ""
        img_b64 = _to_base64(path.read_bytes())
        try:
            result = await llm_client.complete(
                "Transcribe all text from this image. Return only the text.",
                image_data=img_b64,
                image_mime_type=_guess_mime(path),
                image_filename=path.name,
            )
            stats["ocr_called"] = 1
            return result.strip()
        except Exception as exc:
            logger.warning("OCR image %s failed: %s", path.name, exc)
            return ""

    @classmethod
    def _extract_office_text(cls, path: Path, stats: dict) -> str | None:
        """Extract text from Office documents via markitdown or python-docx/openpyxl/pptx."""
        ext = path.suffix.lower()

        # Try python-docx/openpyxl/python-pptx first (faster, no LLM)
        if ext == ".docx":
            text = _extract_docx_fast(path)
            if text: return text
        elif ext == ".xlsx":
            text = _extract_xlsx_fast(path)
            if text: return text
        elif ext == ".pptx":
            text = _extract_pptx_fast(path)
            if text: return text

        # Fall back to markitdown
        return _extract_markitdown(path)

    @classmethod
    def _extract_text_file(cls, path: Path, stats: dict) -> str | None:
        """Read a text file with encoding detection."""
        for encoding in ("utf-8", "gbk", "gb2312", "latin-1"):
            try:
                return path.read_text(encoding=encoding)
            except (UnicodeDecodeError, UnicodeError):
                continue
        return path.read_text(encoding="utf-8", errors="replace")

    @classmethod
    def _extract_unknown(cls, path: Path, stats: dict) -> str | None:
        """Last-resort extraction for unknown file types."""
        try:
            return path.read_text(encoding="utf-8", errors="replace")
        except Exception:
            return ""

    # ── Exam pipeline ────────────────────────────────────────────

    @classmethod
    async def _run_exam_pipeline(
        cls, path: Path, llm_client, max_pages: int, stats: dict,
    ) -> str | None:
        """Run Phase 1-4 structured exam pipeline on a PDF."""
        from tutor_platform.rag.layout_engine import PaperLayoutEngine
        from tutor_platform.rag.block_ocr import BlockOCREngine
        from tutor_platform.rag.exam_structurer import ExamStructurer
        import json as json_mod

        # Phase 1: Layout
        layout = PaperLayoutEngine.process(path)

        # Phase 2: Block OCR
        ocr = BlockOCREngine(llm_client, path)
        try:
            # Limit pages
            pages_to_process = layout.pages[:min(len(layout.pages), max_pages)]
            page_contents = []
            for pg in pages_to_process:
                contents = await ocr.process_page(pg, save_figures=False)
                page_contents.append(contents)
        finally:
            ocr.close()

        # Phase 3: Structure
        exam = ExamStructurer.structure(
            pages_to_process, page_contents, file_hash=layout.file_hash,
        )
        stats["exam_questions"] = len(exam.questions)

        # Phase 4: Serialize
        serialized = {
            "paper_id": exam.paper_id,
            "raw_file_hash": exam.raw_file_hash,
            "total_pages": exam.total_pages,
            "metadata": {
                "subject": exam.metadata.subject,
                "grade": exam.metadata.grade,
                "exam_type": exam.metadata.exam_type,
                "year": exam.metadata.year,
                "total_score": exam.metadata.total_score,
                "duration_minutes": exam.metadata.duration_minutes,
            },
            "questions": [
                {
                    "question_id": q.question_id,
                    "index": q.index,
                    "type": q.type,
                    "content": q.content,
                    "options": q.options,
                    "answer": q.answer,
                    "score": q.score,
                    "page_num": q.page_num,
                    "figures": [
                        {
                            "figure_id": f.figure_id,
                            "bbox": list(f.bbox),
                            "description": f.description,
                        }
                        for f in q.figures
                    ],
                }
                for q in exam.questions
            ],
        }
        return json_mod.dumps(serialized, ensure_ascii=False, indent=2)

    # ── Sidecar helpers ──────────────────────────────────────────

    @classmethod
    def _write_sidecar(cls, file_path: Path, content: str, suffix: str) -> str | None:
        """Write a sidecar text file atomically."""
        sidecar = file_path.with_name(file_path.stem + suffix)
        tmp = sidecar.with_name(sidecar.name + ".tmp")
        try:
            tmp.write_text(content, encoding="utf-8")
            os.replace(str(tmp), str(sidecar))
            return str(sidecar)
        except OSError as exc:
            logger.warning("Failed to write sidecar %s: %s", sidecar, exc)
            return None

    @classmethod
    def _write_sidecar_raw(cls, file_path: Path, content: str, suffix: str) -> str | None:
        """Write a raw sidecar file (not text-transformed)."""
        return cls._write_sidecar(file_path, content, suffix)


# ── Standalone extraction helpers (no LLM) ────────────────────────

def _extract_text_pdf_fast(path: Path) -> str:
    """Extract text from a PDF via pymupdf (fast, no LLM)."""
    try:
        import fitz
    except ImportError:
        return ""
    try:
        doc = fitz.open(path)
        pages = []
        for page in doc:
            pages.append(page.get_text().strip())
        doc.close()
        return "\n\n".join(pages)
    except Exception:
        return ""


def _extract_docx_fast(path: Path) -> str:
    """Extract text from .docx via python-docx."""
    try:
        from docx import Document
    except ImportError:
        return ""
    try:
        doc = Document(str(path))
        return "\n".join(p.text for p in doc.paragraphs if p.text)
    except Exception:
        return ""


def _extract_xlsx_fast(path: Path) -> str:
    """Extract text from .xlsx via openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
        rows: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            for row in ws.iter_rows(values_only=True):
                row_text = "\t".join(str(c) if c is not None else "" for c in row)
                if row_text.strip():
                    rows.append(row_text)
        wb.close()
        return "\n".join(rows)
    except Exception:
        return ""


def _extract_pptx_fast(path: Path) -> str:
    """Extract text from .pptx via python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        prs = Presentation(str(path))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            slide_text.append(p.text.strip())
            if slide_text:
                slides.append(f"Slide {i+1}:\n" + "\n".join(slide_text))
        return "\n\n".join(slides)
    except Exception:
        return ""


def _extract_markitdown(path: Path) -> str:
    """Extract text via Microsoft markitdown."""
    try:
        from markitdown import MarkItDown
    except ImportError:
        return ""
    try:
        md = MarkItDown()
        result = md.convert(str(path))
        return result.text_content if result else ""
    except Exception:
        return ""


def _to_base64(data: bytes) -> str:
    """Encode bytes to base64 string."""
    import base64
    return base64.b64encode(data).decode("ascii")


def _guess_mime(path: Path) -> str:
    """Guess MIME type from extension."""
    ext = path.suffix.lower()
    return {
        ".png": "image/png",
        ".jpg": "image/jpeg",
        ".jpeg": "image/jpeg",
        ".webp": "image/webp",
        ".gif": "image/gif",
        ".bmp": "image/bmp",
    }.get(ext, "application/octet-stream")
