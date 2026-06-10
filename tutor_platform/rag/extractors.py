"""Unified document extractors — single source of truth for text extraction.

Every ingestion path (Web KB upload, WeChat, MCP tools, CLI, chat attachments,
API calls, LlamaIndex document loader) calls these functions so behaviour is
consistent regardless of entry point.

Principles:
  * pymupdf4llm for all PDFs (text layer, scanned, mixed, exams) — one import.
  * python-docx / openpyxl / python-pptx for modern Office — rich structured output.
  * markitdown for EPUB / audio / YouTube / old Office fallback — wide format support.
  * FileTypeRouter for plain text / source code — multi-encoding chain.
  * No Tesseract — OCR is handled by LLM (Qwen2-VL / rkllama) downstream when needed.
"""

from __future__ import annotations

import asyncio
import logging
import os
import re as _re
from pathlib import Path

logger = logging.getLogger(__name__)


# ── PDF extraction ────────────────────────────────────────────────

def extract_pdf_text(
    file_path: str | Path,
    *,
    ocr_enabled: bool = False,
    ocr_trace_id: str = "",
) -> str:
    """Extract text from any PDF as structured Markdown.

    Uses pymupdf4llm.to_markdown() which provides:
      - # heading hierarchy from font-size analysis
      - **bold** / *italic* inline formatting
      - GFM pipe tables from detected table regions
      - Natural reading-order reconstruction for multi-column pages
      - Auto-detected and excluded page headers / footers
      - Embedded image references (![alt](path))

    When ``ocr_enabled=True``, pymupdf4llm's hybrid OCR pipeline is activated:
    only pages/regions without a text layer are OCR'd via the configured
    OCR backend (Qwen2-VL / rkllama via ``ocr_adapters.py``).  Text-layer
    pages pass through unchanged.  No Tesseract required.

    Falls back to raw pymupdf page.get_text() when pymupdf4llm is unavailable
    or raises an exception.
    """
    path = Path(file_path)

    try:
        import pymupdf4llm

        kwargs: dict = {}
        if ocr_enabled:
            from tutor_platform.rag.ocr_adapters import get_minicpm_ocr_function

            kwargs["use_ocr"] = True
            kwargs["ocr_function"] = get_minicpm_ocr_function(trace_id=ocr_trace_id)
            kwargs["force_ocr"] = True
        else:
            kwargs["use_ocr"] = False

        md = pymupdf4llm.to_markdown(path, **kwargs)
        if md and len(md.strip()) > 0:
            logger.debug(
                "pymupdf4llm extracted %d chars from %s (ocr=%s)",
                len(md), path.name, ocr_enabled,
            )
            return md
    except ImportError:
        logger.warning("pymupdf4llm not installed, falling back to pymupdf raw text")
    except Exception as exc:
        logger.warning(
            "pymupdf4llm.to_markdown failed for %s: %s — falling back to pymupdf raw text",
            path.name, exc,
        )

    # Fallback: raw pymupdf page-by-page extraction
    return _extract_pdf_text_fast(path)


def extract_pdf_json(file_path: str | Path) -> list[dict]:
    """Extract structured layout data from a PDF via pymupdf4llm.to_json().

    Returns a list of page dicts with bounding-box, layout type, font metadata
    per element.  Used as input for PaperLayoutEngine downstream.
    """
    path = Path(file_path)

    try:
        import pymupdf4llm
        data = pymupdf4llm.to_json(path, use_ocr=False)
        if data:
            return data
    except ImportError:
        logger.warning("pymupdf4llm not installed, falling back to raw pymupdf")
    except Exception as exc:
        logger.warning("pymupdf4llm.to_json failed for %s: %s", path.name, exc)

    return []


def has_pdf_text_layer(file_path: str | Path, min_chars: int = 200) -> bool:
    """Quick check whether a PDF has meaningful extractable text.

    Strips ``--- Page N ---`` markers (which accumulate on multi-page
    scanned PDFs and can easily exceed a naive char-count threshold)
    before evaluating.

    Returns True when the total *actual* extracted text across all pages
    exceeds ``min_chars`` — meaning the document has a usable digital text
    layer and does NOT need OCR.
    """
    import re as _re
    text = _extract_pdf_text_fast(str(file_path))
    # Strip page markers that pymupdf page.get_text() prepends even for empty pages
    real = _re.sub(r"^--- Page \d+ ---\s*", "", text, flags=_re.MULTILINE).strip()
    return len(real) > min_chars


def extract_pdf_page_count(file_path: str | Path) -> int:
    """Return the number of pages in a PDF without doing full extraction."""
    try:
        import fitz
        with fitz.open(str(file_path)) as doc:
            return len(doc)
    except Exception:
        return 0


# ── Office extraction ──────────────────────────────────────────────

def extract_docx_text(file_path: str | Path) -> str:
    """Extract structured text from .docx using python-docx."""
    try:
        from docx import Document
    except ImportError:
        return ""
    try:
        doc = Document(str(file_path))
        parts: list[str] = []
        for p in doc.paragraphs:
            if p.text and p.text.strip():
                style = p.style.name if p.style else ""
                if "heading" in (style or "").lower() or "title" in (style or "").lower():
                    # Approximate heading level from font size
                    level = 1
                    for run in p.runs:
                        if run.font.size and run.font.size.pt:
                            if run.font.size.pt >= 18:
                                level = 1
                            elif run.font.size.pt >= 14:
                                level = 2
                            break
                    parts.append(f"{'#' * level} {p.text.strip()}")
                else:
                    parts.append(p.text.strip())

        # Include table content
        for table in doc.tables:
            rows: list[str] = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                rows.append("| " + " | ".join(cells) + " |")
            if rows:
                parts.append("\n\n" + "\n".join(rows))

        return "\n\n".join(parts)
    except Exception:
        return _extract_markitdown(str(file_path))


def extract_xlsx_text(file_path: str | Path) -> str:
    """Extract text from .xlsx as GFM tables using openpyxl."""
    try:
        from openpyxl import load_workbook
    except ImportError:
        return ""
    path = Path(file_path)
    try:
        wb = load_workbook(str(path), read_only=True, data_only=True)
        sheets: list[str] = []
        for sheet_name in wb.sheetnames[:10]:  # cap at 10 sheets
            ws = wb[sheet_name]
            rows: list[str] = []
            for row in ws.iter_rows(values_only=True):
                row_text = "| " + " | ".join(str(c) if c is not None else "" for c in row) + " |"
                if row_text.strip("| "):
                    rows.append(row_text)
            if rows:
                sheets.append(f"## Sheet: {sheet_name}\n\n" + "\n".join(rows))
        wb.close()
        return "\n\n".join(sheets)
    except Exception:
        return _extract_markitdown(str(path))


def extract_pptx_text(file_path: str | Path) -> str:
    """Extract text from .pptx using python-pptx."""
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        prs = Presentation(str(file_path))
        slides: list[str] = []
        for i, slide in enumerate(prs.slides):
            slide_text: list[str] = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        if p.text.strip():
                            slide_text.append(p.text.strip())
            if slide_text:
                slides.append(f"### Slide {i + 1}\n\n" + "\n".join(slide_text))
        return "\n\n".join(slides)
    except Exception:
        return _extract_markitdown(str(file_path))


def extract_text_file(file_path: str | Path) -> str:
    """Read a text file with multi-encoding fallback chain."""
    path = Path(file_path)
    for encoding in ("utf-8", "utf-8-sig", "gbk", "gb2312", "gb18030", "latin-1", "cp1252"):
        try:
            return path.read_text(encoding=encoding)
        except (UnicodeDecodeError, UnicodeError):
            continue
    return path.read_text(encoding="utf-8", errors="replace")


def extract_markitdown(file_path: str | Path) -> str:
    """Extract text via Microsoft markitdown (fallback for EPUB, old Office, etc.)."""
    return _extract_markitdown(str(file_path))


# ── Image extraction ──────────────────────────────────────────────

def extract_image_text(file_path: str | Path, *, trace_id: str = "") -> str:
    """Extract text from an image file via OpenCV preprocessing + LLM OCR.

    Steps:
      1. Read raw bytes
      2. OpenCV preprocess (deskew, denoise, CLAHE, threshold)
      3. OCR via configured provider (Qwen2-VL / rkllama)
      4. If full-image OCR fails, try horizontal segment-based OCR

    This function is sync-safe — it blocks the calling thread while
    waiting for OCR HTTP responses (via asyncio.run() internally).
    """
    import asyncio as _asyncio
    import base64 as _base64

    path = Path(file_path)
    try:
        raw_bytes = path.read_bytes()
    except OSError:
        return ""

    try:
        from tutor_platform.tools.preprocess import preprocess_image_bytes
        processed = preprocess_image_bytes(raw_bytes)
    except Exception:
        processed = raw_bytes

    from tutor_platform.rag.ocr_adapters import _ocr_pixmap_bytes

    # ── Run async OCR inside a temporary event loop ──
    async def _do_ocr() -> str:
        # Try full-image OCR first
        text = await _ocr_pixmap_bytes(processed, trace_id)
        if text:
            return text

        # Full-image failed → split into horizontal segments and OCR in parallel
        segments = _split_image_segments(processed)
        if len(segments) <= 1:
            return text

        logger.debug(
            "[%s] Splitting image into %d segments for parallel OCR",
            trace_id, len(segments),
        )
        tasks = [_ocr_pixmap_bytes(seg, trace_id) for seg in segments]
        results = await _asyncio.gather(*tasks)
        combined = "\n".join(r for r in results if r)
        if combined:
            logger.debug(
                "[%s] Segment OCR: %d chars from %d/%d segments",
                trace_id, len(combined), sum(1 for r in results if r), len(results),
            )
            return combined
        return text

    try:
        loop = _asyncio.get_running_loop()
    except RuntimeError:
        return _asyncio.run(_do_ocr())

    # Running loop exists — delegate to a thread
    import concurrent.futures
    with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
        return pool.submit(_asyncio.run, _do_ocr).result()


def _split_image_segments(image_bytes: bytes) -> list[bytes]:
    """Split a page image at horizontal whitespace gaps.

    Uses horizontal projection to detect text row gaps (blank lines between
    questions/paragraphs). Returns the original image as single-element list
    when no split is possible.
    """
    try:
        import cv2
        import numpy as np
    except ImportError:
        return [image_bytes]

    nparr = np.frombuffer(image_bytes, np.uint8)
    img = cv2.imdecode(nparr, cv2.IMREAD_GRAYSCALE)
    if img is None:
        return [image_bytes]

    h, w = img.shape
    _, binary = cv2.threshold(img, 0, 255, cv2.THRESH_BINARY_INV + cv2.THRESH_OTSU)
    proj = np.sum(binary, axis=1) // 255
    threshold_px = max(1, int(w * 0.01))
    content = proj > threshold_px

    min_gap_h = max(5, int(h * 0.005))
    segments: list[bytes] = []
    i = 0
    while i < h:
        if not content[i]:
            i += 1
            continue
        start = i
        while i < h and content[i]:
            i += 1
        end = i
        gap_end = i
        while gap_end < h and not content[gap_end]:
            gap_end += 1
        if gap_end - i >= min_gap_h and end - start >= 40:
            seg = img[max(0, start - 2): min(h, gap_end + 2)]
            _, buf = cv2.imencode(".jpg", seg, [cv2.IMWRITE_JPEG_QUALITY, 90])
            segments.append(buf.tobytes())
            i = gap_end

    if len(segments) <= 1:
        return [image_bytes]
    return segments


# ── PDF table extraction ──────────────────────────────────────────

def extract_pdf_tables(file_path: str | Path, max_pages: int = 50) -> list[dict]:
    """Extract structured table data from a PDF using PyMuPDF ``find_tables()``.

    Returns a list of dicts, each representing one table:
      - ``page``: 0-based page index
      - ``bbox``: table bounding box as ``[x0, y0, x1, y1]``
      - ``row_count``, ``col_count``: table dimensions
      - ``header``: list of column header strings (if detectable)
      - ``rows``: list of lists of cell text

    No Tesseract dependency — tables are detected from text/graphics layout.
    """
    path = Path(file_path)
    try:
        import fitz
    except ImportError:
        logger.warning("PyMuPDF not available for table extraction")
        return []

    tables: list[dict] = []
    try:
        doc = fitz.open(str(path))
        pages_to_check = min(max_pages, len(doc))
        for i in range(pages_to_check):
            page = doc[i]
            found = page.find_tables()
            if not found or not found.tables:
                continue
            for tab in found.tables:
                rows_data: list[list[str]] = []
                header_row: list[str] = []
                for row in tab.extract():
                    cells = [str(cell).strip() if cell is not None else "" for cell in row]
                    if cells and all(cells):
                        rows_data.append(cells)
                if not rows_data:
                    continue
                # First non-empty row is likely the header
                header_row = rows_data[0] if rows_data else []
                data_rows = rows_data[1:] if len(rows_data) > 1 else []
                tables.append({
                    "page": i,
                    "bbox": list(tab.bbox) if hasattr(tab, "bbox") else [],
                    "row_count": len(rows_data),
                    "col_count": len(header_row) if header_row else 0,
                    "header": header_row,
                    "rows": data_rows,
                })
        doc.close()
    except Exception as exc:
        logger.warning("PDF table extraction failed for %s: %s", path.name, exc)

    if tables:
        logger.debug("extract_pdf_tables: %d tables from %s", len(tables), path.name)
    return tables


def extract_pdf_tables_as_markdown(file_path: str | Path, max_pages: int = 50) -> str:
    """Extract PDF tables as GFM Markdown string for appending to extracted text."""
    tables = extract_pdf_tables(file_path, max_pages=max_pages)
    if not tables:
        return ""

    lines = ["\n\n## Tables\n"]
    for i, tab in enumerate(tables):
        lines.append(f"\n### Table {i + 1} (page {tab['page'] + 1})\n")
        header = tab["header"]
        rows = tab["rows"]
        if header:
            lines.append("| " + " | ".join(header) + " |")
            lines.append("|" + "|".join(["---"] * len(header)) + "|")
        for row in rows:
            lines.append("| " + " | ".join(row) + " |")
    return "\n".join(lines)


# ── PDF embedded image extraction ─────────────────────────────────

def extract_pdf_embedded_images(
    file_path: str | Path,
    max_pages: int = 50,
    page_offset: int = 0,
) -> list[bytes]:
    """Extract embedded images from a PDF using ``doc.get_page_images()``.

    Skips the first ``page_offset`` pages (typically front matter: cover,
    copyright, ToC) so callers don't waste the image budget on decorations.

    Returns raw image bytes for each embedded image found.  Caller can
    pass these to an OCR engine for text extraction from embedded graphics.
    """
    path = Path(file_path)
    try:
        import fitz
    except ImportError:
        return []

    images: list[bytes] = []
    try:
        doc = fitz.open(str(path))
        total = len(doc)
        start = min(page_offset, total - 1)
        end = min(start + max_pages, total)
        for i in range(start, end):
            page = doc[i]
            image_list = page.get_images(full=True)
            for img in image_list:
                xref = img[0]
                base_image = doc.extract_image(xref)
                image_bytes = base_image.get("image")
                if image_bytes and len(image_bytes) > 500:
                    images.append(image_bytes)
        doc.close()
    except Exception as exc:
        logger.warning("PDF image extraction failed for %s: %s", path.name, exc)

    if images:
        logger.debug(
            "extract_pdf_embedded_images: %d images from pages %d-%d of %s",
            len(images), start + 1, end, path.name,
        )
    return images


# ── HTML sanitization ─────────────────────────────────────────────

def _sanitize_html(html_text: str) -> str:
    """Strip dangerous HTML elements keeping only safe text content.

    Removes scripts, event handlers, iframes, objects, and style blocks
    to prevent XSS when content is displayed in any web context.
    """
    import re as _re

    cleaned = _re.sub(r"<script[^>]*>.*?</script>", "", html_text,
                      flags=_re.DOTALL | _re.IGNORECASE)
    cleaned = _re.sub(r"<style[^>]*>.*?</style>", "", cleaned,
                      flags=_re.DOTALL | _re.IGNORECASE)
    cleaned = _re.sub(r"<iframe[^>]*>.*?</iframe>", "", cleaned,
                      flags=_re.DOTALL | _re.IGNORECASE)
    cleaned = _re.sub(r"<object[^>]*>.*?</object>", "", cleaned,
                      flags=_re.DOTALL | _re.IGNORECASE)
    cleaned = _re.sub(r"<embed[^>]*>.*?</embed>", "", cleaned,
                      flags=_re.DOTALL | _re.IGNORECASE)
    cleaned = _re.sub(r'\s+on\w+\s*=\s*["\'][^"\']*["\']', "", cleaned,
                      flags=_re.IGNORECASE)
    cleaned = _re.sub(r"\s+on\w+\s*=\s*\S+", "", cleaned, flags=_re.IGNORECASE)
    return cleaned


# ── Unified dispatch ───────────────────────────────────────────────

# Extension → best extractor mapping.  The canonical source of truth.
# All ingestion paths MUST call extract_text() to guarantee consistent behaviour.
_EXTRACTOR_MAP: dict[str, str] = {
    # PDF — pymupdf4llm → pymupdf fallback (OCR via ocr_enabled param)
    ".pdf":  "pdf",
    # Modern Office — dedicated libraries → markitdown fallback
    ".docx": "docx",
    ".xlsx": "xlsx",
    ".pptx": "pptx",
    ".pptm": "pptx",
    ".ppsx": "pptx",
    # Old Office — markitdown only
    ".doc":  "markitdown",
    ".ppt":  "markitdown",
    ".pps":  "markitdown",
    ".xls":  "markitdown",
    ".odt":  "markitdown",
    ".rtf":  "markitdown",
    # Images — OpenCV preprocess + LLM OCR
    ".jpg":  "image",
    ".jpeg": "image",
    ".png":  "image",
    ".gif":  "image",
    ".webp": "image",
    ".bmp":  "image",
    ".tiff": "image",
    ".tif":  "image",
    ".heic": "image",
    # Markitdown formats
    ".epub": "markitdown",
    ".zip":  "markitdown",
    ".mp3":  "markitdown",
    ".wav":  "markitdown",
    ".ogg":  "markitdown",
    # HTML — needs sanitization
    ".html": "html",
    ".htm":  "html",
}


def extract_text(file_path: str | Path, *, trace_id: str = "") -> str:
    """Extract text from any supported document as structured Markdown.

    Automatically selects the best extraction tool based on file extension:
      - .pdf  → pymupdf4llm.to_markdown() (structured Markdown, no OCR by default)
      - .docx → python-docx (headings + tables)
      - .xlsx → openpyxl (GFM tables per sheet)
      - .pptx → python-pptx (slide-by-slide)
      - .jpg/.png/… → OpenCV preprocess + LLM OCR (Qwen2-VL / rkllama)
      - .html → multi-encoding read + XSS sanitization
      - .txt/.md/.py/… → multi-encoding read
      - .doc/.ppt/.xls/.odt/.rtf/.epub → markitdown fallback
      - Others → best-effort text read

    Returns empty string when no text could be extracted.
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    dispatch = _EXTRACTOR_MAP.get(ext, "text")

    if dispatch == "pdf":
        text = extract_pdf_text(path)
        # Cross-validation: if pymupdf4llm returned nothing, try markitdown
        if not text:
            logger.debug("pymupdf4llm returned empty for %s, trying markitdown fallback", path.name)
            text = extract_markitdown(path)
        return text
    elif dispatch == "docx":
        return extract_docx_text(path)
    elif dispatch == "xlsx":
        return extract_xlsx_text(path)
    elif dispatch == "pptx":
        return extract_pptx_text(path)
    elif dispatch == "image":
        return extract_image_text(path, trace_id=trace_id)
    elif dispatch == "html":
        text = extract_text_file(path)
        return _sanitize_html(text) if text else ""
    elif dispatch == "markitdown":
        return extract_markitdown(path)

    # Fallback: try as plain text first, then markitdown
    text = extract_text_file(path)
    if not text and ext not in {
        ".txt", ".md", ".py", ".js", ".ts", ".json", ".yaml", ".yml", ".csv",
        ".html", ".htm", ".xml", ".css", ".log", ".sh", ".cfg", ".ini", ".env",
        ".toml", ".sql", ".tex", ".java", ".c", ".h", ".cpp", ".go", ".rs",
        ".rb", ".php", ".pl", ".lua", ".r", ".swift", ".kt", ".scala", ".dart",
        ".hs", ".clj", ".ex", ".erl", ".ml", ".fs", ".sol", ".proto", ".tf",
    }:
        text = extract_markitdown(path)
    return text or ""


# ── Internal helpers ───────────────────────────────────────────────

def _extract_pdf_text_fast(file_path: str) -> str:
    """Extract text from a PDF via pymupdf page.get_text() (fast, no LLM)."""
    try:
        import fitz
    except ImportError:
        return ""
    try:
        doc = fitz.open(file_path)
        pages = []
        for page in doc:
            pages.append(f"--- Page {page.number + 1} ---\n{page.get_text().strip()}")
        doc.close()
        return "\n\n".join(pages)
    except Exception:
        return ""


def _extract_markitdown(path_str: str) -> str:
    """Extract text via Microsoft markitdown."""
    # For old binary .doc files, try antiword first (markitdown doesn't handle them).
    ext = os.path.splitext(path_str)[1].lower()
    if ext == ".doc":
        text = _extract_doc_antiword(path_str)
        if text:
            return text
    try:
        from markitdown import MarkItDown
    except ImportError:
        return ""
    try:
        md = MarkItDown()
        result = md.convert(path_str)
        return result.text_content if result else ""
    except Exception:
        return ""

def _extract_doc_antiword(file_path: str) -> str:
    """Extract text from old binary .doc files using antiword."""
    import subprocess as _sp
    try:
        result = _sp.run(
            ["antiword", file_path],
            capture_output=True, text=True, timeout=30,
        )
        if result.returncode == 0 and result.stdout.strip():
            return result.stdout.strip()
    except Exception:
        pass
    return ""


# ── Semantic chunking ───────────────────────────────────────────────

def semantic_chunk(
    text: str,
    chunk_size: int | None = None,
    chunk_overlap: int = 100,
    doc_type: str = "unknown",
    filename: str = "",
    figures: list[dict] | None = None,
) -> list[dict]:
    """Split text into semantic chunks preserving Markdown heading boundaries.

    Strategy (in priority order):
      1. Split at ``##`` / ``###`` / ``####`` Markdown headings (preserve
         knowledge-point boundaries).
      2. If resulting chunk > chunk_size, split at ``\\\\n\\\\n`` paragraphs.
      3. Last resort: hard split at chunk_size.

    Each chunk dict contains:
      - ``text``: the chunk content (with overlap prefix)
      - ``heading_path``: breadcrumb of parent headings
      - ``chunk_index``: 0-based chunk number
      - ``char_count``: length of this chunk's own text (excludes overlap)
      - ``doc_type``, ``filename``: metadata passthrough
      - ``figures``: list of figure dicts referenced by or near this chunk
      - ``figure_ids``: list of figure IDs for quick reference matching

    ``chunk_size`` defaults to ``RAG_CHUNK_SIZE`` env var or 800.
    When ``figures`` is provided, each chunk carries a ``figures`` sub-list
    of figure dicts whose spatial/textual context overlaps with that chunk.
    """
    import re as _re

    if chunk_size is None:
        chunk_size = int(os.environ.get("RAG_CHUNK_SIZE", "800"))

    if len(text) <= chunk_size:
        chunk_figures = _match_figures_to_chunk(text, figures or [])
        return [{
            "text": text,
            "heading_path": "",
            "chunk_index": 0,
            "char_count": len(text),
            "doc_type": doc_type,
            "filename": filename,
            "figures": [_figure_ref(f) for f in chunk_figures],
            "figure_ids": [f["figure_id"] for f in chunk_figures],
        }]

    # Split at Markdown headings (## / ### / ####) — each starts a new section
    heading_re = _re.compile(r"^(#{2,4})\s+(.+)$", _re.MULTILINE)
    lines = text.split("\n")
    sections: list[tuple[str, str]] = []  # [(heading_path, body_text)]
    current_heading = ""
    current_body: list[str] = []

    heading_stack: list[tuple[int, str]] = []  # [(level, title)]

    for line in lines:
        m = heading_re.match(line.strip())
        if m and line.startswith("#"):
            # Flush current section
            if current_body:
                body_text = "\n".join(current_body).strip()
                if body_text:
                    sections.append((current_heading, body_text))
                current_body = []

            level = len(m.group(1))
            title = m.group(2).strip()

            # Pop headings that are same or deeper level
            while heading_stack and heading_stack[-1][0] >= level:
                heading_stack.pop()
            heading_stack.append((level, title))
            current_heading = " > ".join(t for _, t in heading_stack)
        else:
            current_body.append(line)

    # Flush final section
    if current_body:
        body_text = "\n".join(current_body).strip()
        if body_text:
            sections.append((current_heading, body_text))

    # Build chunks from sections, splitting oversized ones
    _figures = figures or []
    chunks: list[dict] = []
    for heading, body in sections:
        chunk_figures = _match_figures_to_chunk(body, _figures)
        if len(body) <= chunk_size:
            chunks.append({
                "text": body,
                "heading_path": heading,
                "chunk_index": len(chunks),
                "char_count": len(body),
                "doc_type": doc_type,
                "filename": filename,
                "figures": [_figure_ref(f) for f in chunk_figures],
                "figure_ids": [f["figure_id"] for f in chunk_figures],
            })
        else:
            # Split oversized body at paragraph boundaries
            paragraphs = _re.split(r"\n\s*\n", body)
            buf = ""
            buf_heading = heading
            for para in paragraphs:
                if len(buf) + len(para) + 2 > chunk_size and buf:
                    buf_text = buf.strip()
                    pfigures = _match_figures_to_chunk(buf_text, _figures)
                    chunks.append({
                        "text": buf_text,
                        "heading_path": buf_heading,
                        "chunk_index": len(chunks),
                        "char_count": len(buf_text),
                        "doc_type": doc_type,
                        "filename": filename,
                        "figures": [_figure_ref(f) for f in pfigures],
                        "figure_ids": [f["figure_id"] for f in pfigures],
                    })
                    buf = para
                else:
                    buf = (buf + "\n\n" + para).strip() if buf else para
            if buf.strip():
                buf_text = buf.strip()
                pfigures = _match_figures_to_chunk(buf_text, _figures)
                chunks.append({
                    "text": buf_text,
                    "heading_path": heading,
                    "chunk_index": len(chunks),
                    "char_count": len(buf_text),
                    "doc_type": doc_type,
                    "filename": filename,
                    "figures": [_figure_ref(f) for f in pfigures],
                    "figure_ids": [f["figure_id"] for f in pfigures],
                })

    # Apply overlap: prepend tail of previous chunk
    if chunk_overlap > 0 and len(chunks) > 1:
        for i in range(1, len(chunks)):
            prev = chunks[i - 1]["text"]
            if len(prev) > chunk_overlap:
                overlap = prev[-chunk_overlap:]
                # Find a clean break point
                nl = overlap.find("\n")
                if nl > chunk_overlap // 2:
                    overlap = overlap[nl + 1:]
                chunks[i]["text"] = overlap + "\n" + chunks[i]["text"]

    return chunks


# ── Figure-chunk matching helpers ──────────────────────────────────

_RE_FIGURE_REF = _re.compile(
    r"[图图表示][0-9一二三四五六七八九十\-]|"
    r"如图|Fig\.\s*\d+|figure\s*\d+|"
    r"所示|示意图|图\d+",
    _re.IGNORECASE,
)


def _match_figures_to_chunk(chunk_text: str, figures: list[dict]) -> list[dict]:
    """Match figures to a chunk by checking for textual references.

    A figure is considered "referenced" by a chunk when the chunk text
    contains ``如图X`` / ``Fig.X`` / ``图X`` patterns, or when the
    figure's caption/OCR text overlaps with the chunk content.

    Returns a subset of ``figures`` that are relevant to this chunk.
    """
    if not figures or not chunk_text:
        return []

    matched: list[dict] = []
    seen_ids: set[str] = set()

    # Check for explicit figure references in the chunk text
    has_ref = bool(_RE_FIGURE_REF.search(chunk_text))

    for fig in figures:
        fid = fig.get("figure_id", "")
        if fid in seen_ids:
            continue

        # Match by textual reference pattern
        caption = fig.get("caption", "")
        fig_type = fig.get("fig_type", "")
        ocr_text = fig.get("ocr_text", "")

        # If the chunk mentions this figure's caption or OCR text
        if caption and caption in chunk_text:
            seen_ids.add(fid)
            matched.append(fig)
            continue

        if has_ref and (caption or fig_type != "unknown"):
            seen_ids.add(fid)
            matched.append(fig)
            continue

        # Last-resort: if OCR text from the figure appears verbatim in chunk
        if ocr_text and len(ocr_text) > 10 and ocr_text in chunk_text:
            seen_ids.add(fid)
            matched.append(fig)
            continue

    return matched


def _figure_ref(fig_dict: dict) -> dict:
    """Return a compact, JSON-safe reference to a figure for embedding in chunk metadata."""
    return {
        "figure_id": fig_dict.get("figure_id", ""),
        "fig_type": fig_dict.get("fig_type", "unknown"),
        "caption": fig_dict.get("caption", ""),
        "description_text": fig_dict.get("description_text", ""),
        "page_num": fig_dict.get("page_num", 0),
    }


# ── Unified figure extraction ─────────────────────────────────────

_IMAGE_EXTENSIONS: set[str] = {
    ".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp", ".tiff", ".tif", ".heic",
}


def extract_figures(
    file_path: str | Path,
    *,
    trace_id: str = "",
    llm_client=None,
    max_figures: int = 60,
) -> list[dict]:
    """Extract all figures from a document as serialisable dicts.

    Dispatches to the appropriate sub-extractor based on file extension.
    Returns a list of ``UnifiedFigure``-compatible dicts (JSON-safe).

    Args:
        file_path: Path to the source document.
        trace_id: Trace ID for logging.
        llm_client: Optional multimodal LLM client for figure description.
                    When ``None``, only OCR text is extracted (no description).
        max_figures: Maximum number of figures to extract (capped to avoid
                     excessive NPU/LLM usage on large documents).
    """
    path = Path(file_path)
    ext = path.suffix.lower()

    if ext == ".pdf":
        return _extract_pdf_figures(path, trace_id=trace_id, llm_client=llm_client,
                                     max_figures=max_figures)
    elif ext in _IMAGE_EXTENSIONS:
        return _extract_image_figure(path, trace_id=trace_id, llm_client=llm_client)
    elif ext in (
        ".docx", ".pptx", ".pptm", ".ppsx", ".xlsx",
        ".doc", ".ppt", ".pps", ".xls",
    ):
        return _extract_office_figures(path, trace_id=trace_id, llm_client=llm_client,
                                        max_figures=max_figures)
    return []


def _extract_pdf_figures(
    path: Path,
    *,
    trace_id: str = "",
    llm_client=None,
    max_figures: int = 20,
) -> list[dict]:
    """Extract figures from a PDF via ``extract_pdf_embedded_images()`` +
    optional multimodal VL structured description.

    If a ``.exam.json`` sidecar already exists next to the PDF, its
    ``QuestionFigure`` entries are reused instead of re-running the LLM.
    """
    from tutor_platform.rag.figure_types import UnifiedFigure

    # ── Prefer existing .exam.json sidecar ──
    sidecar = path.with_name(path.stem + ".exam.json")
    if sidecar.exists():
        try:
            return _consume_exam_sidecar(sidecar, source_file=str(path))
        except Exception as exc:
            logger.warning("[%s] Failed to consume exam sidecar %s: %s",
                           trace_id, sidecar.name, exc)

    # ── Fallback: extract embedded images directly ──
    # Skip first 15 pages (cover, copyright, ToC), then scan 120 pages
    # of content, capped at 60 figures per PDF.
    _cap = min(max_figures, 60)
    raw_images = extract_pdf_embedded_images(path, max_pages=120, page_offset=15)
    if not raw_images:
        return []

    figures: list[dict] = []
    for idx, img_bytes in enumerate(raw_images[:_cap]):
        fig = UnifiedFigure(
            source_file=str(path),
            page_num=0,
            image_bytes=img_bytes,
            fig_type="unknown",
        )

        # OCR text in the figure if an LLM client is available
        if llm_client and hasattr(llm_client, "complete"):
            try:
                import asyncio, base64
                img_b64 = base64.b64encode(img_bytes).decode("ascii")
                # Reuse the text OCR prompt from the exam pipeline
                result = _run_llm_ocr_sync(llm_client, img_b64)
                if result:
                    fig.ocr_text = result
                    fig.fig_type = _infer_fig_type_from_text(result)
            except Exception as exc:
                logger.warning("[%s] LLM OCR failed for PDF figure %d: %s",
                               trace_id, idx, exc)
        else:
            # ── RapidOCR fallback (no LLM client available) ──
            try:
                from tutor_platform.rag.rapid_ocr import ocr_text_only
                ocr_result = ocr_text_only(img_bytes)
                if ocr_result:
                    fig.ocr_text = ocr_result
                    fig.fig_type = _infer_fig_type_from_text(ocr_result)
            except ImportError:
                pass
            except Exception as exc:
                logger.debug("[%s] RapidOCR figure fallback failed: %s", trace_id, exc)

        figures.append(_figure_to_dict(fig))

    if figures:
        logger.debug("[%s] Extracted %d figures from %s",
                      trace_id, len(figures), path.name)
    return figures


def _extract_image_figure(
    path: Path,
    *,
    trace_id: str = "",
    llm_client=None,
) -> list[dict]:
    """Extract a figure from a standalone image file.

    Combines OCR text extraction and LLM vision description into a single
    ``UnifiedFigure``.
    """
    from tutor_platform.rag.figure_types import UnifiedFigure

    raw_bytes = path.read_bytes()
    fig = UnifiedFigure(
        source_file=str(path),
        page_num=0,
        image_bytes=raw_bytes,
        fig_type="illustration",
    )

    # OCR text first
    try:
        ocr_text = extract_image_text(path, trace_id=trace_id)
        if ocr_text:
            fig.ocr_text = ocr_text
    except Exception as exc:
        logger.warning("[%s] Image OCR failed for %s: %s", trace_id, path.name, exc)

    # Vision description (separate call, cached externally)
    if llm_client and hasattr(llm_client, "complete"):
        try:
            import asyncio, base64
            img_b64 = base64.b64encode(raw_bytes).decode("ascii")
            prompt = (
                "Describe this image. Include visible text/OCR if present, "
                "the main subject, and any educational or technical meaning. "
                "Keep under 180 words."
            )
            result = _run_llm_ocr_sync(llm_client, img_b64, prompt)
            if result:
                fig.description = {"raw": result}
        except Exception as exc:
            logger.warning("[%s] Vision description failed for %s: %s",
                           trace_id, path.name, exc)

    # If no LLM, use OCR text directly
    if not fig.description and fig.ocr_text:
        fig.description = {"raw": f"Image containing text: {fig.ocr_text[:200]}"}

    result = _figure_to_dict(fig)
    logger.debug("[%s] Extracted figure from %s (ocr=%d, desc=%d)",
                  trace_id, path.name,
                  len(fig.ocr_text or ""),
                  len(str(fig.description or "")))
    return [result]


def _extract_office_figures(
    path: Path,
    *,
    trace_id: str = "",
    llm_client=None,
    max_figures: int = 20,
) -> list[dict]:
    """Extract figures from Office documents (docx/pptx/xlsx).

    Upgrades the existing ``_extract_office_images()`` listing to include
    OCR text and LLM description for each embedded image.
    """
    from tutor_platform.rag.figure_types import UnifiedFigure

    raw_images = _extract_office_images_for_figures(path)
    if not raw_images:
        return []

    figures: list[dict] = []
    for idx, img_bytes in enumerate(raw_images[:max_figures]):
        fig = UnifiedFigure(
            source_file=str(path),
            page_num=0,
            image_bytes=img_bytes,
            fig_type="illustration",
        )

        if llm_client and hasattr(llm_client, "complete"):
            try:
                import base64
                img_b64 = base64.b64encode(img_bytes).decode("ascii")
                result = _run_llm_ocr_sync(llm_client, img_b64)
                if result:
                    fig.ocr_text = result
            except Exception as exc:
                logger.warning("[%s] Office figure OCR failed for %s img %d: %s",
                               trace_id, path.name, idx, exc)
        else:
            # ── RapidOCR fallback (no LLM client available) ──
            try:
                from tutor_platform.rag.rapid_ocr import ocr_text_only
                ocr_result = ocr_text_only(img_bytes)
                if ocr_result:
                    fig.ocr_text = ocr_result
            except ImportError:
                pass
            except Exception as exc:
                logger.debug("[%s] RapidOCR figure fallback failed: %s", trace_id, exc)

        figures.append(_figure_to_dict(fig))

    if figures:
        logger.debug("[%s] Extracted %d figures from Office %s",
                      trace_id, len(figures), path.name)
    return figures


def _extract_office_images_for_figures(path: Path) -> list[bytes]:
    """Extract all embedded images from an Office document.

    Enhanced version of ``_extract_office_images()`` that returns raw
    image bytes for LLM processing instead of just a text listing.

    Handles both ZIP-based formats (docx/pptx/xlsx) and OLE-based
    old formats (doc/ppt/xls).
    """
    import zipfile, os as _os

    ext = path.suffix.lower()
    all_images: list[bytes] = []
    seen_hashes: set[int] = set()

    # ── ZIP-based formats (modern Office) ──
    media_prefixes = {
        ".docx": ["word/media/"], ".docm": ["word/media/"],
        ".pptx": ["ppt/media/"], ".pptm": ["ppt/media/"], ".ppsx": ["ppt/media/"],
        ".xlsx": ["xl/media/"],
    }.get(ext, [])

    for prefix in media_prefixes:
        try:
            with zipfile.ZipFile(str(path)) as z:
                for name in z.namelist():
                    ext_lower = _os.path.splitext(name)[1].lower()
                    if name.startswith(prefix) and ext_lower in {
                        ".png", ".jpg", ".jpeg", ".gif", ".bmp",
                    }:
                        data = z.read(name)
                        h = hash(data)
                        if h not in seen_hashes and len(data) > 500:
                            seen_hashes.add(h)
                            all_images.append(data)
        except (zipfile.BadZipFile, FileNotFoundError):
            pass

    # ── OLE-based formats (old Office: .doc, .ppt, .xls) ──
    if ext in {".doc", ".ppt", ".pps", ".xls"}:
        try:
            import olefile
        except ImportError:
            return all_images

        try:
            ole = olefile.OleFileIO(str(path))
            for s in ole.listdir():
                try:
                    data = ole.openstream(s).read()
                    h = hash(data)
                    if h in seen_hashes:
                        continue
                    sig = data[:8]
                    if (
                        sig[:2] == b"\xff\xd8"
                        or sig[:4] == b"\x89PNG"
                        or sig[:2] == b"BM"
                        or sig[:4] == b"GIF8"
                    ):
                        seen_hashes.add(h)
                        all_images.append(data)
                except Exception:
                    continue
            ole.close()
        except Exception:
            pass

    return all_images


def _consume_exam_sidecar(sidecar_path: Path, source_file: str = "") -> list[dict]:
    """Parse an ``.exam.json`` sidecar and return UnifiedFigure dicts."""
    import json as _json
    from tutor_platform.rag.figure_types import UnifiedFigure

    with open(sidecar_path, "r", encoding="utf-8") as f:
        exam = _json.load(f)

    figures: list[dict] = []
    for question in exam.get("questions", []):
        for qf in question.get("figures", []):
            desc = qf.get("description")
            fig_type = _infer_fig_type_from_desc(desc) if desc else "unknown"
            fig = UnifiedFigure(
                figure_id=qf.get("figure_id", ""),
                source_file=source_file or str(sidecar_path),
                page_num=qf.get("page_num", 0),
                bbox=tuple(qf.get("bbox", [0, 0, 0, 0])),
                description=desc,
                fig_type=fig_type,
            )
            # If the description has OCR-like text, pull it out
            if desc and isinstance(desc, dict):
                for key in ("text", "ocr", "content"):
                    val = desc.get(key)
                    if val and isinstance(val, str):
                        fig.ocr_text = val
                        break
            figures.append(_figure_to_dict(fig))

    return figures


# ── Helpers ────────────────────────────────────────────────────────

def _figure_to_dict(fig) -> dict:
    """Convert a UnifiedFigure to a JSON-serialisable dict."""
    return {
        "figure_id": fig.figure_id,
        "source_file": fig.source_file,
        "page_num": fig.page_num,
        "bbox": list(fig.bbox) if fig.bbox else None,
        "image_bytes": fig.image_bytes,
        "image_path": fig.image_path,
        "ocr_text": fig.ocr_text,
        "description": fig.description,
        "fig_type": fig.fig_type,
        "caption": fig.caption,
        "referring_chunks": fig.referring_chunks,
        "description_text": fig.description_text,
    }


def _infer_fig_type_from_text(text: str) -> str:
    """Guess figure type from OCR text content."""
    import re as _re
    tl = text.lower()
    if any(w in tl for w in ("三角形", "circle", "圆", "矩形", "正方形",
                              "平行", "垂直", "∠", "△", "▱", "□")):
        return "geometry"
    if _re.search(r'[=xysincostanlim]', tl) and _re.search(r'\d', tl):
        return "function_graph"
    if _re.search(r'[│|]\s*[^│|]+\s*[│|]', text) or "table" in tl:
        return "table"
    if any(w in tl for w in ("图", "diagram", "chart", "实验", "示意")):
        return "illustration"
    return "unknown"


def _infer_fig_type_from_desc(desc: dict) -> str:
    """Determine figure type from a structured description dict."""
    if not isinstance(desc, dict):
        return "unknown"
    if "figure_type" in desc:
        return "geometry"
    if "function_hint" in desc:
        return "function_graph"
    if "type" in desc:
        t = str(desc.get("type", "")).lower()
        if t in ("table",):
            return "table"
    return "illustration"


def _run_llm_ocr_sync(llm_client, img_b64: str, prompt: str = "") -> str:
    """Run a synchronous LLM OCR/description call.

    Works with both sync and async LLM clients by detecting the running
    event loop.
    """
    import asyncio as _asyncio

    if not prompt:
        prompt = (
            "Transcribe all visible text from this image exactly as written, "
            "preserving the original language and structure. "
            "If the image contains diagrams or figures, describe them "
            "briefly. Return only the transcribed text and description."
        )

    async def _call() -> str:
        return await llm_client.complete(
            prompt,
            image_data=img_b64,
            image_mime_type="image/png",
            image_filename="figure.png",
        )

    try:
        loop = _asyncio.get_running_loop()
        if loop.is_running():
            import concurrent.futures
            with concurrent.futures.ThreadPoolExecutor(max_workers=1) as pool:
                return pool.submit(_asyncio.run, _call()).result()
        return _asyncio.run(_call())
    except RuntimeError:
        return _asyncio.run(_call())


def _match_answers(student: str, correct: str) -> bool:
    """Compare student answer with correct answer, handles various formats.

    Supports:
    - Exact match (选择题: "B" == "B")
    - Case-insensitive match
    - Numeric equivalence ("33" == "33岁", "x=5" == "5")
    - Trimmed comparison (trailing units/punctuation)
    """
    s = student.strip()
    c = correct.strip()
    if not s or not c:
        return False

    if s == c:
        return True
    if s.upper() == c.upper():
        return True
    if s in ("A", "B", "C", "D") and c in ("A", "B", "C", "D"):
        return False
    s_nums = _re.findall(r"\d+", s)
    c_nums = _re.findall(r"\d+", c)
    if s_nums and c_nums and all(cn in s_nums for cn in c_nums):
        return True
    s_clean = _re.sub(r"[.。，,、\s单位个只条约根种]+$", "", s)
    c_clean = _re.sub(r"[.。，,、\s单位个只条约根种]+$", "", c)
    if s_clean and c_clean and s_clean == c_clean:
        return True
    return False


def _match_answers_semantic(student: str, correct: str) -> bool:
    """Partial match detection — student has the right idea but not exact."""
    import re as _re2
    s = student.strip().lower()
    c = correct.strip().lower()
    if not s or not c or s == c:
        return False

    _opt_s = _re2.findall(r"^选?([a-dA-D])$|\(?([a-dA-D])\)?$", s)
    _opt_c = _re2.findall(r"^选?([a-dA-D])$|\(?([a-dA-D])\)?$", c)
    if _opt_s and _opt_c:
        _s_letter = (_opt_s[0][0] or _opt_s[0][1]).upper()
        _c_letter = (_opt_c[0][0] or _opt_c[0][1]).upper()
        if _s_letter == _c_letter:
            return True

    _s_nums = set(_re2.findall(r"\d+", s))
    _c_nums = set(_re2.findall(r"\d+", c))
    if _s_nums and _c_nums:
        if len(_s_nums & _c_nums) >= len(_c_nums) * 0.5:
            return True

    _key_terms = {
        "加", "减", "乘", "除", "等于", "大", "小", "正", "负",
        "数", "和", "差", "积", "商", "解", "方程", "分式",
        "分子", "分母", "倒数", "绝对", "平方", "根",
    }
    _s_terms = {t for t in _key_terms if t in s}
    _c_terms = {t for t in _key_terms if t in c}
    if _s_terms and _c_terms and _s_terms == _c_terms:
        return True
    return False
